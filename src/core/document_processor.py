"""
Document Processor Module
Handles PDF (digital + scanned), Word, PowerPoint with OCR
Optimized for low RAM usage
"""
import os
import re
from pathlib import Path
from typing import List, Dict

try:
    import pytesseract
    from PIL import Image
    PYTESSERACT_AVAILABLE = True
except ImportError:
    PYTESSERACT_AVAILABLE = False
    pytesseract = None

try:
    from pdf2image import convert_from_path
    PDF2IMAGE_AVAILABLE = True
except ImportError:
    PDF2IMAGE_AVAILABLE = False
    convert_from_path = None

import PyPDF2
import pdfplumber
from docx import Document
from pptx import Presentation

import config


class DocumentProcessor:
    """Process various document types and extract text"""
    
    def __init__(self):
        # Set Tesseract path if specified in environment
        if PYTESSERACT_AVAILABLE:
            tesseract_path = os.getenv('TESSERACT_CMD')
            if not tesseract_path:
                # Default Windows installation path
                tesseract_path = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
            
            if os.path.exists(tesseract_path):
                pytesseract.pytesseract.tesseract_cmd = tesseract_path
            else:
                print(f"⚠ Tesseract not found at {tesseract_path}")

        # Check if Tesseract is available
        self.tesseract_available = self._check_tesseract()
        self.pdf2image_available = PDF2IMAGE_AVAILABLE

    def _check_tesseract(self) -> bool:
        """Check if Tesseract is available"""
        if not PYTESSERACT_AVAILABLE:
            return False
        try:
            pytesseract.get_tesseract_version()
            return True
        except Exception:
            return False
    
    def process_document(self, file_path: str) -> List[str]:
        """
        Process document and return list of text chunks
        
        Args:
            file_path: Path to document file
            
        Returns:
            List of text chunks (≤500 tokens each)
        """
        file_path = Path(file_path)
        
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        # Extract text based on file type
        extension = file_path.suffix.lower()
        
        if extension == '.pdf':
            text = self._process_pdf(file_path)
        elif extension == '.docx':
            text = self._process_docx(file_path)
        elif extension == '.pptx':
            text = self._process_pptx(file_path)
        elif extension == '.txt':
            text = self._process_txt(file_path)
        else:
            raise ValueError(f"Unsupported file type: {extension}")
        
        # Clean and chunk the text
        cleaned_text = self._clean_text(text)
        chunks = self._chunk_text(cleaned_text)
        
        return chunks
    
    def _process_pdf(self, file_path: Path) -> str:
        """
        Process PDF file (digital or scanned)
        Try digital extraction first, fall back to OCR
        """
        text = ""
        
        # Try digital extraction first using pdfplumber (better than PyPDF2)
        try:
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
        except Exception as e:
            print(f"pdfplumber failed: {e}, trying PyPDF2...")
        
        # If pdfplumber didn't extract much, try PyPDF2
        if len(text.strip()) < 50:
            try:
                with open(file_path, 'rb') as file:
                    reader = PyPDF2.PdfReader(file)
                    for page in reader.pages:
                        page_text = page.extract_text()
                        if page_text:
                            text += page_text + "\n"
            except Exception as e:
                print(f"PyPDF2 failed: {e}")
        
        # If still not much text, it's probably scanned - use OCR
        if len(text.strip()) < 50:
            print(f"Little/no text extracted ({len(text.strip())} chars). Attempting OCR...")
            if self.tesseract_available and self.pdf2image_available:
                try:
                    text = self._ocr_pdf(file_path)
                    print(f"OCR extracted {len(text)} characters")
                except Exception as e:
                    print(f"OCR failed: {e}")
                    # Try alternative OCR method using pdfplumber images
                    text = self._ocr_pdf_alternative(file_path)
            else:
                missing = []
                if not self.tesseract_available:
                    missing.append("Tesseract")
                if not self.pdf2image_available:
                    missing.append("pdf2image")
                print(f"⚠ Scanned PDF detected but {', '.join(missing)} not available.")
                # Try alternative OCR with pdfplumber
                text = self._ocr_pdf_alternative(file_path)
        
        if not text.strip():
            raise ValueError("No text could be extracted from this PDF. It may be empty, corrupted, or image-only without proper OCR support.")
        
        return text
    
    def _ocr_pdf(self, file_path: Path) -> str:
        """Perform OCR on scanned PDF"""
        text = ""
        
        try:
            # Set poppler path for Windows
            poppler_path = r'C:\poppler\poppler-24.08.0\Library\bin'
            
            # Convert PDF pages to images (process one page at a time for low RAM)
            pages = convert_from_path(
                file_path,
                dpi=200,  # Lower DPI for RAM efficiency
                fmt='jpeg',
                grayscale=True,  # Reduce memory usage
                poppler_path=poppler_path if os.path.exists(poppler_path) else None
            )
            
            # OCR each page
            for i, page_image in enumerate(pages):
                print(f"OCR processing page {i+1}/{len(pages)}...")
                page_text = pytesseract.image_to_string(
                    page_image,
                    lang=config.OCR_LANGUAGE
                )
                text += page_text + "\n"
                
                # Free memory
                page_image.close()
                del page_image
            
        except Exception as e:
            print(f"OCR (pdf2image) failed: {e}")
            raise
        
        return text
    
    def _ocr_pdf_alternative(self, file_path: Path) -> str:
        """
        Alternative OCR method using pdfplumber image extraction
        Works even without pdf2image
        """
        text = ""
        
        if not PYTESSERACT_AVAILABLE:
            print("⚠ Cannot perform OCR: pytesseract not available")
            return text
        
        try:
            from PIL import Image
            import io
            
            with pdfplumber.open(file_path) as pdf:
                for page_num, page in enumerate(pdf.pages):
                    print(f"OCR processing page {page_num+1}/{len(pdf.pages)} (alternative method)...")
                    
                    # Try to extract images from the page
                    images = page.images
                    
                    if images:
                        # If page has images, OCR them
                        for img in images:
                            try:
                                # This is a simplified approach
                                # In practice, pdfplumber doesn't easily give us the image data
                                # So we'll render the whole page as image
                                pass
                            except Exception as e:
                                print(f"Image extraction failed: {e}")
                    
                    # Alternative: render the whole page as image and OCR
                    # This requires page.to_image() which needs Pillow
                    try:
                        img = page.to_image(resolution=150)
                        # Convert to PIL Image
                        page_text = pytesseract.image_to_string(img.original)
                        text += page_text + "\\n"
                    except Exception as e:
                        print(f"Page rendering failed: {e}")
                        
        except Exception as e:
            print(f"Alternative OCR failed: {e}")
        
        return text
    
    def _process_docx(self, file_path: Path) -> str:
        """Process Word document (.docx)"""
        text = ""
        
        try:
            doc = Document(file_path)
            
            # Extract text from paragraphs
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text += cell.text + " "
                text += "\n"
                
        except Exception as e:
            print(f"Error processing DOCX: {e}")
            raise
        
        return text
    
    def _process_pptx(self, file_path: Path) -> str:
        """Process PowerPoint document (.pptx)"""
        text = ""
        
        try:
            prs = Presentation(file_path)
            
            # Extract text from all slides
            for slide_num, slide in enumerate(prs.slides, 1):
                text += f"\n--- Slide {slide_num} ---\n"
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                    
                    # Extract text from tables in slides
                    if shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                text += cell.text + " "
                        text += "\n"
                        
        except Exception as e:
            print(f"Error processing PPTX: {e}")
            raise
        
        return text
    
    def _process_txt(self, file_path: Path) -> str:
        """Process plain text file (.txt)"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()
        except UnicodeDecodeError:
            # Try with different encoding if UTF-8 fails
            with open(file_path, 'r', encoding='latin-1') as f:
                text = f.read()
        return text
    
    def _clean_text(self, text: str) -> str:
        """Clean extracted text"""
        # Remove excessive whitespace
        text = re.sub(r'\s+', ' ', text)
        
        # Remove special characters but keep punctuation
        text = re.sub(r'[^\w\s.,!?;:\-\(\)\'\"]+', '', text)
        
        # Remove multiple consecutive punctuation
        text = re.sub(r'([.,!?;:])\1+', r'\1', text)
        
        return text.strip()
    
    def _chunk_text(self, text: str, max_tokens: int = None) -> List[str]:
        """
        Split text into chunks of maximum token size
        
        Args:
            text: Input text
            max_tokens: Maximum tokens per chunk (default from config)
            
        Returns:
            List of text chunks
        """
        if max_tokens is None:
            max_tokens = config.CHUNK_SIZE
        
        # Approximate: 1 token ≈ 4 characters (conservative estimate)
        max_chars = max_tokens * 4
        overlap_chars = config.CHUNK_OVERLAP * 4
        
        chunks = []
        sentences = re.split(r'(?<=[.!?])\s+', text)
        
        current_chunk = ""
        
        for sentence in sentences:
            # If adding this sentence exceeds max, save current chunk
            if len(current_chunk) + len(sentence) > max_chars and current_chunk:
                chunks.append(current_chunk.strip())
                
                # Start new chunk with overlap
                words = current_chunk.split()
                overlap_words = words[-int(overlap_chars/5):] if len(words) > 10 else []
                current_chunk = " ".join(overlap_words) + " " + sentence
            else:
                current_chunk += " " + sentence
        
        # Add the last chunk
        if current_chunk.strip():
            chunks.append(current_chunk.strip())
        
        return chunks
    
    def process_directory(self, directory_path: str) -> Dict[str, List[str]]:
        """
        Process all documents in a directory
        
        Args:
            directory_path: Path to directory containing documents
            
        Returns:
            Dictionary mapping filename to list of chunks
        """
        directory = Path(directory_path)
        supported_extensions = ['.pdf', '.docx', '.pptx', '.txt']
        
        results = {}
        
        for file_path in directory.iterdir():
            if file_path.suffix.lower() in supported_extensions:
                print(f"\nProcessing: {file_path.name}")
                try:
                    chunks = self.process_document(str(file_path))
                    results[file_path.name] = chunks
                    print(f"[OK] Extracted {len(chunks)} chunks from {file_path.name}")
                except Exception as e:
                    print(f"[ERROR] Failed to process {file_path.name}: {e}")
        
        return results


if __name__ == "__main__":
    # Test the document processor
    processor = DocumentProcessor()
    
    # Process all documents in the documents directory
    docs_dir = config.DOCUMENTS_DIR
    print(f"Processing documents from: {docs_dir}\n")
    
    results = processor.process_directory(str(docs_dir))
    
    print(f"\n{'='*60}")
    print(f"Total documents processed: {len(results)}")
    total_chunks = sum(len(chunks) for chunks in results.values())
    print(f"Total chunks created: {total_chunks}")
    print(f"{'='*60}")
